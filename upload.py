import os
from mimetypes import MimeTypes 
import mimetypes
from pydrive.auth import GoogleAuth
from pydrive.drive import GoogleDrive
import win32com.client
import pythoncom
import glob
import time
import io
import pandas as pd
import docx
from docx import Document
from docx.shared import Inches
import regex
from pptx import Presentation
from pptx.enum.shapes import MSO_SHAPE_TYPE
from nltk.corpus import stopwords
import nltk
from gensim.summarization import keywords
import json
import re
from flask import send_file,url_for
def retrieve(id = 'TEST_DBMS'):
    subject_folder = id.split('_')[-1]
    if os.path.exists('./Upload Folder/' + subject_folder):
        path = './Upload Folder/' + subject_folder + '/' + id
        name_links = os.listdir(path=path)
        print(name_links)
        link_name = []
        for a in name_links:
            download_link = url_for('download',id = id,filename = a)
            link_name.append((download_link,a))
        return link_name
    return []

def auth():
    if os.path.exists('mycreds.txt'):
        gauth = GoogleAuth()
        gauth.LoadCredentialsFile("mycreds.txt")
    else:
        gauth = GoogleAuth()
        gauth.LocalWebserverAuth()
        gauth.SaveCredentialsFile("mycreds.txt")
    drive = GoogleDrive(gauth)
    return drive

def upload(id,to,frompath):
    #drive = auth()
    name = frompath.split('\\')[-1]
    extension = frompath.split('.')[-1]
    if extension == 'pdf':
        path = PDF2WORD(frompath,name,to)
        name = name[:-4]
        name += '.docx'
    elif extension == 'pptx':
        path = PPT2WORD(frompath,name,to)
        name = name[:-5]
        name += '.docx'
    genkeywords(path,'./Upload Folder/' + id.split('_')[-1] + '/',id)
    

def PDF2WORD(path,name,to):
    pythoncom.CoInitialize()
    folder_loc = to
    doc = glob.glob(path)[0]
    word = win32com.client.Dispatch("Word.Application")
    word.Visible = True
    pdfloc = os.path.abspath(doc)
    name = name[:-4]
    wb = word.Documents.Open(pdfloc)
    doc_path = os.path.abspath(folder_loc + "\\" + name + '.docx')
    wb.SaveAs2(doc_path,FileFormat = 16)
    wb.Close()
    word.Quit()
    time.sleep(2)
    return doc_path
def PPT2WORD(path,name,to):
    pres = Presentation(path)
    folder_loc = to
    name = name[:-5]
    doc = Document()
    for slide in pres.slides:
        for shape in slide.shapes:
            if shape == slide.shapes.title:
                t = slide.shapes.title.text
                if t != '' and t.isprintable():
                    p = doc.add_paragraph().add_run(t)
                    p.bold = True
            if shape != slide.shapes.title and (shape.shape_type == MSO_SHAPE_TYPE.PLACEHOLDER or shape.shape_type == MSO_SHAPE_TYPE.TEXT_BOX):
                if hasattr(shape,"text"):
                    para = doc.add_paragraph()
                    text = shape.text
                    text = ''.join(c for c in text if ord(c) >= 32 or c == '\n' or c == '\t')
                    for c in text:
                        if c.isprintable():
                            para.add_run(c)
                        else:
                            s = para.add_run(c)
                            font = s.font
                            font.name = 'Symbol'
            if shape.shape_type == MSO_SHAPE_TYPE.PICTURE or shape.shape_type == MSO_SHAPE_TYPE.PLACEHOLDER:
                if hasattr(shape,"image"):
                    image = shape.image
                    image_bin = image.blob
                    try:
                        doc.add_picture(io.BytesIO(image_bin),width = Inches(4),height = Inches(4))
                        doc.add_paragraph()
                    except:
                        print('Exception')
            if shape.shape_type == MSO_SHAPE_TYPE.TABLE or shape.shape_type == MSO_SHAPE_TYPE.PLACEHOLDER:
                if hasattr(shape,"table"):
                    table = shape.table
                    try:
                        doc_table = doc.add_table(0,len(table.columns))
                        for i in range(0,len(table.rows)):
                            row = doc_table.add_row().cells
                            for j in range(0,len(table.columns)):
                                temp = table.cell(i,j).text
                                row[j].text = temp
                        doc_table.style = 'TableGrid'
                    except:
                        print('Exception')
                    doc.add_paragraph()
    doc.save('{}.docx'.format(folder_loc + "\\" + name))
    time.sleep(2)
    return folder_loc + "\\" + name + ".docx"


def genkeywords(path,path1,id):
    stopword = stopwords.words('english')
    AUTHOR = id
    FILENAME = path
    doc = docx.Document(FILENAME)
    document = doc
    map = {}
    r = []
    for paragraph in document.paragraphs:
        if "Heading" not in paragraph.style.name:
            runboldtext = ''
            normal = ''
            for run in paragraph.runs:                        
                if run.bold:
                    runboldtext = runboldtext + run.text
                else:
                    normal += run.text
            if runboldtext == str(paragraph.text) and runboldtext != '':
                r.append((runboldtext.strip(),True))
                #print("Heading True for the paragraph: ",runboldtext)
                style_of_current_paragraph = 'Heading'
            else:
                #print("Heading False for the paragraph: ",normal)
                r.append((normal,False))
        else:
        #print("Heading True for the paragraph: ",paragraph.text)
            r.append((paragraph.text.strip(),True))
    i = 0
    cur = ''
    while i < len(r):
        if r[i][0] == '':
            i += 1
            continue
        elif r[i][1] == True:
            temp = r[i][0]
            i += 1
            while i < len(r) and r[i][1] == False:
                cur += r[i][0]
                i += 1
            if cur != '':
                if temp in map.keys():
                    map[temp] += cur
                else:
                    map[temp] = cur
                cur = ''
        elif r[i][1] == False:
            map['No heading'] = r[i][0]
            i += 1
    l = []
    for key,val in map.items():
        if len(val.strip()) == 0 or len(val.split()) < 20  or len(val) == 0:
            continue
        imp_words = []
        try:
            imp_words.append(keywords(val,words = 5,split = True,lemmatize = True,scores = False))
        except:
            continue
        if key != 'No Heading':
            imp_words.append([word for word in key.split() if word not in stopword])
        imp_words = [word for words in imp_words for word in words]
        imp_words = preprocess(imp_words)
        name = FILENAME.split('\\')[-1]
        details = {'Author' : AUTHOR, 'Chapter' : name[:-5], 'Topic' : key, 'Keywords' : imp_words, 'Paragraph' : val}
        l.append(details)
    json_file = path1 + 'metadata.json'
    if os.path.exists(json_file):
        with open(json_file,'r') as f:
            temp = json.load(f)
            for ls in l:
                temp.append(ls)
        with open(json_file,'w') as f:
            json.dump(temp,f)
    else:
        with open(json_file,'w') as f:
            json.dump(l,f)
    if os.path.exists(json_file):
        with open(json_file,'r') as f:
            df = json.load(f)
        df = pd.DataFrame.from_dict(df,orient = 'columns')
        map = {}
        for ind,row in enumerate(df['Keywords']):
            for word in row:
                if word not in map.keys():
                    map[word] = [ind]
                else:
                    temp = map[word]
                    temp.append(ind)
                    map[word] = temp
        with open(path1 + 'keyword_indices.json','w') as f:
            json.dump(map,f)
def preprocess(data):
  other_words = []
  stop_words = set(stopwords.words('english')) 
  for ind,row in enumerate(data):
    row = row.lower()
    row = re.sub('\s+',' ',row)
    row = re.sub(r'[^a-zA-z0-9\s]','',row)
    if len(row.split()) == 1:
      row = row.strip()
    else:
      row = ''
      for f in row.split():
        other_words.append(f)
    data[ind] = row
  data = [f for f in data if f != '' and f not in stop_words]
  for f in other_words:
    data.append(f)
  return list(set(data))

if __name__ == '__main__':
    print(retrieve())
